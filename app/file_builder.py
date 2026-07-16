"""Native Office artifact builders — always produce valid binary files for preview."""

import os
import re
from typing import List, Tuple


def _is_zip_magic(path: str) -> bool:
    try:
        with open(path, "rb") as f:
            return f.read(2) == b"PK"
    except OSError:
        return False


def _parse_md_table_lines(lines: List[str], start: int) -> Tuple[list, list, int]:
    """Parse a GFM markdown table starting at `start`. Returns (headers, rows, next_index)."""
    if start + 1 >= len(lines):
        return [], [], start
    header_line = lines[start].strip()
    sep_line = lines[start + 1].strip()
    if not header_line.startswith("|") or "---" not in sep_line:
        return [], [], start

    def split_row(line: str) -> list:
        return [c.strip() for c in line.strip().strip("|").split("|")]

    headers = split_row(header_line)
    rows = []
    i = start + 2
    while i < len(lines):
        row_line = lines[i].strip()
        if not row_line.startswith("|"):
            break
        rows.append(split_row(row_line))
        i += 1
    return headers, rows, i


def build_docx(path: str, content: str) -> None:
    import docx

    doc = docx.Document()
    lines = content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        stripped = line.strip()

        if not stripped:
            doc.add_paragraph()
            i += 1
            continue

        # Markdown table block
        if stripped.startswith("|") and i + 1 < len(lines) and "---" in lines[i + 1]:
            headers, rows, i = _parse_md_table_lines(lines, i)
            if headers:
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = "Table Grid"
                for ci, h in enumerate(headers):
                    cell = table.rows[0].cells[ci]
                    cell.text = h.replace("**", "")
                    for p in cell.paragraphs:
                        for r in p.runs:
                            r.bold = True
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row[: len(headers)]):
                        table.rows[ri + 1].cells[ci].text = val.replace("**", "")
                doc.add_paragraph()
            continue

        if stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            p = doc.add_paragraph(style="List Bullet")
            _add_rich_runs(p, stripped[2:])
        else:
            p = doc.add_paragraph()
            _add_rich_runs(p, stripped)
        i += 1

    doc.save(path)
    if not _is_zip_magic(path):
        raise ValueError("Generated file is not a valid DOCX archive")


def _add_rich_runs(paragraph, text: str) -> None:
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            paragraph.add_run(part[2:-2]).bold = True
        elif part:
            paragraph.add_run(part)


def build_pptx(path: str, content: str) -> None:
    from pptx import Presentation

    prs = Presentation()
    # Split on slide breaks: --- or ## Slide / # Slide
    chunks = re.split(r"\n---+\n", content.strip())
    if len(chunks) == 1:
        chunks = re.split(r"(?=\n##\s+)", content.strip())

    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        lines = chunk.split("\n")
        title = "Slide"
        body_lines = []
        for line in lines:
            s = line.strip()
            if s.startswith("# "):
                title = s[2:]
            elif s.startswith("## "):
                title = s[3:]
            elif s:
                body_lines.append(s.lstrip("- ").replace("**", ""))

        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        slide.shapes.title.text = title
        if body_lines and len(slide.placeholders) > 1:
            slide.placeholders[1].text = "\n".join(body_lines[:12])

    if not prs.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Document"
        slide.placeholders[1].text = content[:500]

    prs.save(path)
    if not _is_zip_magic(path):
        raise ValueError("Generated file is not a valid PPTX archive")


def build_xlsx(path: str, content: str) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    lines = [l for l in content.strip().split("\n") if l.strip()]

    if lines and lines[0].strip().startswith("|") and len(lines) > 1 and "---" in lines[1]:
        headers, rows, _ = _parse_md_table_lines(lines, 0)
        for ci, h in enumerate(headers, 1):
            ws.cell(row=1, column=ci, value=h.replace("**", ""))
        for ri, row in enumerate(rows, 2):
            for ci, val in enumerate(row[: len(headers)], 1):
                ws.cell(row=ri, column=ci, value=val.replace("**", ""))
    else:
        for ri, line in enumerate(lines, 1):
            cells = [c.strip() for c in line.split(",")]
            for ci, val in enumerate(cells, 1):
                ws.cell(row=ri, column=ci, value=val)

    wb.save(path)
    if not _is_zip_magic(path):
        raise ValueError("Generated file is not a valid XLSX archive")


def extract_pptx_slides(path: str) -> dict:
    """Extract slide text for frontend PptxPreview."""
    from pptx import Presentation

    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        title = ""
        body = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if not title:
                title = text
            else:
                body.extend(text.split("\n"))
        slides.append({"title": title, "body": body})
    return {"slides": slides}


def write_artifact_file(path: str, content: str) -> str:
    """Write a workspace artifact; returns success or error message."""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            build_docx(path, content)
        elif ext == ".pptx":
            build_pptx(path, content)
        elif ext == ".xlsx":
            build_xlsx(path, content)
        elif ext in (".doc", ".ppt", ".xls"):
            return (
                f"Error: Legacy '{ext}' is not supported. "
                f"Use '{ext}x' instead (e.g. report{ext}x)."
            )
        else:
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)

        if ext in (".docx", ".pptx", ".xlsx") and not os.path.getsize(path):
            return f"Error: Failed to write non-empty {ext} file."

        # Persist to Firestore so the file survives server restarts
        try:
            import base64
            from app.storage import save_json
            filename = os.path.basename(path)
            if ext in (".docx", ".pptx", ".xlsx", ".pdf", ".png", ".jpg", ".jpeg", ".zip"):
                with open(path, "rb") as f:
                    encoded = base64.b64encode(f.read()).decode("utf-8")
                save_json("heccker_artifact_files", filename, {"filename": filename, "encoding": "base64", "data": encoded}, "")
            else:
                save_json("heccker_artifact_files", filename, {"filename": filename, "encoding": "utf-8", "data": content}, "")
        except Exception:
            pass  # Don't fail the write just because Firestore is unavailable

        return f"Success: Content written to '{path}'."
    except ImportError as e:
        return f"Error: Missing library for {ext} — {e}"
    except Exception as e:
        if os.path.exists(path):
            try:
                os.remove(path)
            except OSError:
                pass
        return f"Error writing {ext}: {e}"
